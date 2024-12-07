import os
import cv2
import traceback
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

from datetime import datetime




def create_pdf_from_frames(frames, video_filename, output_folder, logger=None):
    """Enhanced PDF creation with timeline-based frame titles and robust error handling."""
    try:
        output_pdf = os.path.join(output_folder, f"{os.path.splitext(video_filename)[0]}.pdf")
        
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        for i in range(0, len(frames), 3):
            pdf.add_page()
            y_positions = [10, 90, 170]
            
            for j, frame_data in enumerate(frames[i:i+3]):
                frame_path, timestamp = frame_data
                
                try:
                    # Add image to PDF
                    pdf.image(frame_path, x=10, y=y_positions[j], w=190, h=80)
                    pdf.set_y(y_positions[j] + 80)

                    # Add frame and timestamp as caption
                    pdf.set_font("Arial", size=10)
                    pdf.cell(0, 10, f"Frame: {os.path.basename(frame_path)} : {timestamp}", ln=True, align="C")

                except Exception as img_error:
                    if logger:
                        logger.warning(f"Could not add frame {frame_path}: {img_error}")
        
        pdf.output(output_pdf)
        
        if logger:
            logger.info(f"PDF created: {output_pdf}")
        
        return output_pdf

    except Exception as e:
        if logger:
            logger.error(f"PDF creation error: {e}")
            logger.error(traceback.format_exc())
        return None



def create_pptx_from_frames(frames, video_filename, output_folder, logger=None):
    """Enhanced PPTX creation with timeline logging."""
    try:
        output_pptx = os.path.join(output_folder, f"{os.path.splitext(video_filename)[0]}.pptx")
        prs = Presentation()

        for idx, frame_data in enumerate(frames):
            frame_path, timestamp = frame_data
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            frame_img = cv2.imread(frame_path)
            if frame_img is None:
                if logger:
                    logger.warning(f"Could not read frame {frame_path}")
                continue

            frame_width, frame_height = frame_img.shape[1], frame_img.shape[0]

            # Intelligent image scaling
            if frame_width / slide_width > frame_height / slide_height:
                width = slide_width
                height = int(frame_height * (slide_width / frame_width))
            else:
                height = slide_height
                width = int(frame_width * (slide_height / frame_height))

            left = (slide_width - width) // 2
            top = (slide_height - height) // 2

            slide.shapes.add_picture(frame_path, left, top, width=width, height=height)

            # Add slide title with frame and timestamp
            title = slide.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(4), Inches(0.5))
            title_frame = title.text_frame
            title_frame.text = f"Frame {idx + 1} : {timestamp}"
            title_frame.paragraphs[0].font.size = Pt(18)

        prs.save(output_pptx)

        if logger:
            logger.info(f"PPTX created: {output_pptx}")

        return output_pptx

    except Exception as e:
        if logger:
            logger.error(f"PPTX creation error: {e}")
            logger.error(traceback.format_exc())
        return None


def extract_frames(video_path, interval, output_dir="frames", logger=None):
    """
    Enhanced frame extraction with timeline logging.

    Args:
        video_path (str): Path to the input video file
        interval (int): Seconds between frame extractions
        output_dir (str): Directory to save extracted frames
        logger (logging.Logger): Logger for tracking extraction process

    Returns:
        list: List of tuples (frame_path, timestamp)
    """
    try:

        if logger:
            logger.info(f"Starting frame extraction from video: {video_path}")
            logger.info(f"Frame extraction interval: {interval} seconds")
            
            # Adding human-readable units for file size
            file_size = os.path.getsize(video_path)
            human_readable_size = f"{file_size / (1024 ** 2):.2f} MB" if file_size >= 1024 ** 2 else f"{file_size / 1024:.2f} KB"
            logger.info(f"Video file size: {human_readable_size}")
            
            # Formatting modification time as a readable date and time
            modification_time = os.path.getmtime(video_path)
            readable_time = datetime.fromtimestamp(modification_time).strftime("%Y-%m-%d %H:%M:%S")
            logger.info(f"Last modified on: {readable_time}")


        os.makedirs(output_dir, exist_ok=True)

        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            if logger:
                logger.error(f"Failed to open video: {video_path}")
            return []

        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        frame_interval = max(1, int(fps * interval))

        extracted_frames = []
        frame_count = 0

        if logger:
            logger.info(f"Processing video: {video_path}")
            logger.info(f"Total Frames: {total_frames}, FPS: {fps}, Extraction Interval: {frame_interval}")

        while frame_count < total_frames:
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_count)
            ret, frame = cap.read()

            if not ret:
                break

            # Calculate timestamp
            current_time_sec = frame_count / fps
            timestamp = f"{int(current_time_sec // 3600):02}:{int((current_time_sec % 3600) // 60):02}:{int(current_time_sec % 60):02}"

            # Save frame
            frame_name = os.path.join(output_dir, f"frame_{frame_count:05d}.jpg")
            cv2.imwrite(frame_name, frame, [cv2.IMWRITE_JPEG_QUALITY, 85])
            extracted_frames.append((frame_name, timestamp))

            # if logger:
            #     logger.info(f"Extracted Frame {frame_count // frame_interval + 1} at {timestamp}")

            frame_count += frame_interval

        cap.release()

        if logger:
            logger.info(f"Extracted {len(extracted_frames)} frames from {video_path}")

        return extracted_frames

    except Exception as e:
        if logger:
            logger.error(f"Frame extraction error for {video_path}: {e}")
            logger.error(traceback.format_exc())
        return []


